"""
Generate aa-to-cja-migration-component-manager.pptx from dim28_Master.potx
"""
import shutil, zipfile, os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

TEMPLATE = r"C:\Users\lukas\PycharmProjects\adobe-analytics-component-manager-doc\dim28_Master.potx"
OUTPUT   = r"C:\Users\lukas\PycharmProjects\adobe-analytics-component-manager-doc\aa-to-cja-migration-component-manager.pptx"

# ── patch .potx → loadable as .pptx ─────────────────────────────────────────
tmp = TEMPLATE.replace(".potx", "_tmp.pptx")
shutil.copy(TEMPLATE, tmp)
with zipfile.ZipFile(tmp, 'r') as z:
    contents = {n: z.read(n) for n in z.namelist()}
ct = contents['[Content_Types].xml'].decode()
ct = ct.replace(
    'presentationml.template.main+xml',
    'presentationml.presentation.main+xml'
)
contents['[Content_Types].xml'] = ct.encode()
with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as z:
    for name, data in contents.items():
        z.writestr(name, data)

prs = Presentation(tmp)
os.remove(tmp)

# ── layout indices ───────────────────────────────────────────────────────────
L_TITLE       = 0   # Titelfolie
L_CHAPTER     = 2   # Kapiteltrenner
L_DARK_1      = 3   # Titel mit Inhalt, dunkel
L_DARK_2      = 4   # Titel mit 2 Inhalten, dunkel
L_DARK_2_KEY  = 8   # Titel mit 2 Inhalten und Kernaussage
L_CLOSING     = 24  # Abschlussfolie

# ── helpers ──────────────────────────────────────────────────────────────────
def add_slide(layout_idx):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def ph(slide, idx):
    for p in slide.placeholders:
        if p.placeholder_format.idx == idx:
            return p
    return None

def set_title(slide, text):
    p = ph(slide, 0)
    if p:
        p.text = text

def set_subtitle(slide, text):
    p = ph(slide, 1)
    if p:
        p.text = text

def fill(slide, ph_idx, lines):
    """Fill a content placeholder with bullet lines.
    Prefix with '  ' for level-1 indent, '    ' for level-2.
    Empty string '' adds a small spacer paragraph.
    """
    placeholder = ph(slide, ph_idx)
    if not placeholder:
        return
    tf = placeholder.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        if line == "":
            para.text = ""
            continue
        level = 0
        text = line
        if line.startswith("    "):
            level = 2
            text = line.strip()
        elif line.startswith("  "):
            level = 1
            text = line.strip()
        para.level = level
        run = para.add_run()
        run.text = text

# ── SLIDE 1 — Title ──────────────────────────────────────────────────────────
s = add_slide(L_TITLE)
set_title(s, "Prepare for the Migration from\nAdobe Analytics to Customer Journey Analytics")
set_subtitle(s, "Using the Component Manager  ·  datacroft.de")

# ── SLIDE 2 — A Failed Migration ─────────────────────────────────────────────
s = add_slide(L_DARK_1)
set_title(s, "A Failed Migration")
fill(s, 12, [
    '"Where is my dashboard?"',
    '"Why are the numbers so different?"',
    '"Where is that segment I used in AA?"',
    "",
    "  Frustrated early adopters spread bad word-of-mouth instantly",
    "  CJA's near-identical interface creates false expectations",
    "  Dimensions missing, segments gone, reports rebuilt from scratch",
    "",
    "The Component Manager helps prevent this.",
    "",
    "You can never make a second first impression.",
])

# ── SLIDE 3 — Why CJA Migrations Disappoint (2-col) ─────────────────────────
s = add_slide(L_DARK_2)
set_title(s, "Why CJA Migrations Often Disappoint")
fill(s, 12, [
    "Looks the same…",
    "",
    "  Interface nearly identical to AA",
    "  Same concepts: segments, calc metrics,",
    "  dimensions, workspaces",
    "  Power users feel at home immediately",
    "",
    "  → Sets very high expectations",
])
fill(s, 13, [
    "…but works very differently",
    "",
    "  No eVars/props → data model changes",
    "  Segments don't transfer 1:1",
    "  VRS logic is different",
    "  Thousands of AA components",
    "  won't map cleanly",
    "",
    "  → Unmet expectations = project failure",
])

# ── SLIDE 4 — What is the Component Manager? ────────────────────────────────
s = add_slide(L_DARK_1)
set_title(s, "What is the Component Manager?")
fill(s, 12, [
    "A Google Sheets Add-on for Adobe Analytics Admins",
    "",
    "  Full overview of all AA components in one sheet:",
    "  Segments · Calc Metrics · eVars/props · Success Events · VRSs",
    "",
    "  Bulk-edit, bulk-delete, replace components across Workspaces",
    "  Component Usage: see what's used — and where",
    "  Account Usage: track user adoption, identify power users",
    "",
    "For migration: it shows you what to keep and what to leave behind.",
])

# ── SLIDE 5 — Can't we just use Adobe's Component Migration Tool? ────────────
s = add_slide(L_DARK_2)
set_title(s, "Can't We Just Use Adobe's Component Migration Tool?")
fill(s, 12, [
    "What it does:",
    "  Recreates AA Workspaces, segments, calc metrics",
    "  and date ranges in CJA",
    "",
    "The limitations:",
    "  Needs a 1:1 match of AA components → XDM fields",
    "  (3 eVars in AA often = 1 XDM field in CJA)",
    "  Some AA logic simply won't work in CJA",
    "  (e.g. visitor-profile-based segment conditions)",
    "  Some panel types not supported at all",
    "  (Analytics for Target, Page Summary, …)",
])
fill(s, 13, [
    "The deeper problem:",
    "  Chicken-and-egg: needs near-full CJA config",
    "  to map to — but you need to know what",
    "  you need in CJA first",
    "",
    "  Cements old AA thinking (eVars/props)",
    "  into the new system",
    "",
    '  "We avoided the Migration Tool and built',
    '  a new, streamlined setup instead."',
    "  — Patrick Hegnauer, Adobe Champion",
    "",
    "  The tool still doesn't tell you WHAT",
    "  to migrate. That's where the",
    "  Component Manager comes in.",
])

# ── SLIDE 6 — Chapter: The Strategy ─────────────────────────────────────────
s = add_slide(L_CHAPTER)
set_title(s, "The Strategy")
set_subtitle(s, "Four steps to a leaner CJA setup and a smooth migration")

# ── SLIDE 6 — Strategy Overview ─────────────────────────────────────────────
s = add_slide(L_DARK_1)
set_title(s, "Overview: Four Phases")
fill(s, 12, [
    "Part 1 — Clean up AA before you touch CJA",
    "  Delete dead Workspaces & components. Harmonize duplicates.",
    "",
    "Part 2 — Identify what's actually used in the Workspaces worth migrating",
    "  Limit Component Usage to relevant Workspaces only.",
    "",
    "Part 3 — Rethink your data model (eVars / Success Events)",
    "  The migration is your once-in-a-decade chance to reduce complexity.",
    "",
    "Part 4 — Win your users — power users as allies, not spectators",
    "",
    "Everything you remove now = lower migration cost, leaner CJA.",
])

# ── SLIDE 7 — Part 1: Clean Up ───────────────────────────────────────────────
s = add_slide(L_DARK_1)
set_title(s, "Part 1: Clean Up First")
fill(s, 12, [
    "Step 1 — Identify and delete dead Workspaces",
    "  Use Workspace view stats to find what nobody has opened in months",
    "",
    "Step 2 — Identify and delete dead components",
    "  Unused segments, calc metrics, date ranges visible in Component Usage tab",
    "",
    "Step 3 — Harmonize duplicates with the Component Replacer",
    "  Some segments exist under 10 names with identical definitions",
    "  Replace all variants with the canonical one — across all Workspaces",
    "",
    "Every item deleted now = one less item to migrate, document & maintain in CJA.",
])

# ── SLIDE 8 — Part 2: Identify What's Worth Migrating ───────────────────────
s = add_slide(L_DARK_2)
set_title(s, "Part 2: Identify What's Worth Migrating")
fill(s, 12, [
    "Even after cleanup: still too many",
    "components to migrate everything.",
    "",
    "Filter Component Usage to only the",
    "'relevant' Workspaces:",
    "",
    "  Step 4: Copy Workspaces tab",
    "  → e.g. wsp_to_limit_comp_usage",
    "",
    "  Step 5: Filter to Workspaces worth",
    "  migrating (views, tags, manual picks)",
    "  Tip: email owners to flag priorities",
])
fill(s, 13, [
    "  Step 6: Set limit_workspaces_for_",
    "  comp_usage in config tab",
    "",
    "  Step 7: Set run_with_hard_refresh",
    "  = TRUE",
    "",
    "  Step 8: Run Component Usage Update",
    "",
    "matching_projects_count now shows",
    "usage only within relevant Workspaces.",
    "",
    "→ These components are your",
    "  CJA migration candidates.",
])

# ── SLIDE 9 — Part 3: Rethink the Data Model ────────────────────────────────
s = add_slide(L_DARK_1)
set_title(s, "Part 3: Rethink the Data Model")
fill(s, 12, [
    "AA eVars/props/Success Events ≠ CJA dimensions/metrics",
    "— but usage patterns transfer directly.",
    "",
    "  High usage → must be in the CJA tracking plan",
    "  Low usage  → strong case for leaving it out of CJA entirely",
    "",
    "Step 9: Check dimension & metric usage in Component Usage tab",
    "  Low usage: exclude from CJA tracking plan",
    "  High usage: confirm it's covered in CJA",
    "",
    "Step 10 (optional): In rs_editor, find eVars/events with no recent data",
    "  → Exclude 'dead' variables from CJA setup",
    "",
    "The migration is your once-in-a-decade chance to truly reduce complexity.",
])

# ── SLIDE 10 — Part 4: Win Your Users (2-col) ───────────────────────────────
s = add_slide(L_DARK_2)
set_title(s, "Part 4: Win Your Users")
fill(s, 12, [
    "Step 11: POC with Power Users",
    "",
    "  Account Usage tab → find your most",
    "  active AA users",
    "",
    "  Don't surprise them with a big-bang",
    "  rollout",
    "",
    "  Run a POC migration for their key",
    "  Workspaces together",
    "",
    "  Make them advocates, not skeptics.",
])
fill(s, 13, [
    "Step 12: Support the Stragglers",
    "",
    "  As CJA becomes the default, use",
    "  Account Usage to find users still",
    "  active in AA",
    "",
    "  Why are they still there?",
    "    Missing feature?",
    "    Need targeted CJA onboarding?",
    "",
    "  Ensure a clean AA end-of-life —",
    "  nobody left behind.",
])

# ── SLIDE 12 — Key Takeaways ─────────────────────────────────────────────────
s = add_slide(L_DARK_2_KEY)
set_title(s, "Key Takeaways")
fill(s, 12, [
    "Migration is architecture,",
    "not relocation.",
    "",
    "  Whatever you don't clean up before",
    "  the migration, you'll carry with you",
    "  for the next ten years.",
    "",
    "Adobe's Component Migration Tool",
    "needs a selection first —",
    "the selection is the actual work.",
    "",
    "  The Component Manager helps you make",
    "  that selection in a data-driven,",
    "  user-centric manner.",
])
fill(s, 13, [
    "Power users decide whether",
    "the migration succeeds.",
    "",
    "  Lose them early, and the migration",
    "  is politically dead.",
    "",
    "  Involve them as POC partners —",
    "  not as spectators.",
])
key = ph(s, 14)
if key:
    key.text = "Don't migrate what nobody uses."

# ── SLIDE 12 — Questions / Closing ──────────────────────────────────────────
s = add_slide(L_CLOSING)
set_title(s, "Comments & Questions?")
set_subtitle(s, "datacroft.de  ·  component-manager@datacroft.de")

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
